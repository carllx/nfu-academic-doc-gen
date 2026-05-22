#!/usr/bin/env python3
"""
NFU (广州南方学院) PPT 品牌注入脚本

工作模式："固定环节 + 内容合并"
  - 固定环节：封面、回顾、目录、引用、作业、结尾（由本脚本生成）
  - 内容主体：由用户已有的课程 PPT 提供
  - 可选环节：章节分隔页

最终输出：[封面] → [回顾] → [目录] → [原有内容...] → [引用] → [作业] → [结尾]

使用方式：
  python inject_branding.py --input content.pptx --output branded.pptx \\
      --course-name "交互设计" --teacher "张三" ...

依赖：python-pptx, pyyaml, lxml
"""

import argparse
import os
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.oxml.ns import qn
    from lxml import etree
except ImportError:
    print("错误：缺少依赖。请执行：pip install python-pptx lxml")
    sys.exit(1)

try:
    import yaml
except ImportError:
    print("错误：缺少 PyYAML。请执行：pip install pyyaml")
    sys.exit(1)

# ====================================================================
# 路径解析（确保跨项目可迁移—全部相对于脚本自身位置）
# ====================================================================
SKILL_DIR = Path(__file__).resolve().parent.parent
RESOURCES_DIR = SKILL_DIR / "resources"
THEME_PATH = RESOURCES_DIR / "nfu_theme.yaml"


def load_theme() -> dict:
    """加载 NFU 设计 Token"""
    with open(THEME_PATH, "r", encoding="utf-8") as f:
        return yaml.safe_load(f)


def resolve_media_dir(aspect: str) -> Path:
    """根据比例选择媒体素材目录"""
    return RESOURCES_DIR / "media" / ("16-9" if aspect == "16:9" else "4-3")


def detect_aspect(prs: Presentation) -> str:
    """探测幻灯片宽高比"""
    return "16:9" if (prs.slide_width / prs.slide_height) > 1.5 else "4:3"


# ====================================================================
# OOXML 低级工具
# ====================================================================

def _hex_rgb(hex_str: str) -> RGBColor:
    return RGBColor.from_string(hex_str)


def _find_blank_layout(prs: Presentation):
    """定位空白布局（优先 layoutType='blank' / 第 7 个 / 最后一个）"""
    for layout in prs.slide_layouts:
        el = layout.element
        if el.get("type") == "blank":
            return layout
    if len(prs.slide_layouts) > 6:
        return prs.slide_layouts[6]
    return prs.slide_layouts[-1]


def _move_slide(prs: Presentation, from_idx: int, to_idx: int):
    """将幻灯片从 from_idx 移动到 to_idx"""
    sldIdLst = prs.element.find(qn("p:sldIdLst"))
    entries = list(sldIdLst)
    target = entries[from_idx]
    sldIdLst.remove(target)
    if to_idx == 0:
        sldIdLst.insert(0, target)
    else:
        list(sldIdLst)[to_idx - 1].addnext(target)


def _set_ea_font(run, font_name: str):
    """为 run 设置东亚字体"""
    rPr = run._r.get_or_add_rPr()
    ea = rPr.find(qn("a:ea"))
    if ea is None:
        ea = etree.SubElement(rPr, qn("a:ea"))
    ea.set("typeface", font_name)


def _set_line_spacing(paragraph, pct: int):
    """设置段落行距百分比"""
    pPr = paragraph._p.get_or_add_pPr()
    lnSpc = pPr.find(qn("a:lnSpc"))
    if lnSpc is None:
        lnSpc = etree.SubElement(pPr, qn("a:lnSpc"))
    for child in list(lnSpc):
        lnSpc.remove(child)
    spcPct = etree.SubElement(lnSpc, qn("a:spcPct"))
    spcPct.set("val", str(pct * 1000))


# ====================================================================
# 通用文本构建器
# ====================================================================

def _add_fullscreen_bg(slide, image_path: str, w: int, h: int):
    """添加全屏背景图片"""
    slide.shapes.add_picture(image_path, 0, 0, w, h)


def _add_logo(slide, logo_path: str, cfg: dict):
    """在 layout 配置指定位置放置 Logo"""
    lc = cfg["logo"]
    slide.shapes.add_picture(logo_path, lc["x"], lc["y"], lc["cx"], lc["cy"])


def _add_textbox(slide, pos: dict, lines: list, theme: dict,
                 first_line_sz_key: str = "cover_title_sz",
                 body_sz_key: str = "meta_sz",
                 color_hex: str = "FFFFFF",
                 bold: bool = False,
                 spacing_pct: int = 160):
    """通用文本框构建器"""
    typo = theme["typography"]
    cn_font = typo["cn_primary"]

    txBox = slide.shapes.add_textbox(pos["x"], pos["y"], pos["cx"], pos["cy"])
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, text in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT

        run = p.add_run()
        run.text = text

        # 字号：首行用大标题，其余用正文
        sz_key = first_line_sz_key if i == 0 else body_sz_key
        if sz_key in typo:
            run.font.size = Pt(typo[sz_key] / 100)

        run.font.name = cn_font
        run.font.bold = bold
        run.font.color.rgb = _hex_rgb(color_hex)
        _set_ea_font(run, cn_font)
        _set_line_spacing(p, spacing_pct)

    return txBox


# ====================================================================
# 固定环节页生成函数
# ====================================================================

def gen_cover(prs, theme, layout, media, info, *, position=0):
    """封面页：深灰背景 + 课程名/编号/教师/学期 + Logo"""
    slide = prs.slides.add_slide(_find_blank_layout(prs))
    w, h = layout["slide_width"], layout["slide_height"]
    _add_fullscreen_bg(slide, str(media / theme["media"]["cover_bg"]["file"]), w, h)
    _add_logo(slide, str(media / theme["media"]["logo_combo"]["file"]), layout)

    lines = [
        info.get("course_name", "课程名称"),
        f"课程编号：{info.get('course_code', '')}",
        "",
        f"主讲教师：{info.get('teacher', '')}",
        f"开课单位：{theme['brand']['name_cn']}",
        f"授课学期：{info.get('semester', '')}",
        f"授课时间：{info.get('date', '')}",
    ]
    _add_textbox(slide, layout["cover_text"], lines, theme,
                 color_hex=theme["palette"]["text_on_dark"],
                 spacing_pct=160)

    if position == 0:
        _move_slide(prs, len(prs.slides) - 1, 0)
    return slide


def gen_review(prs, theme, layout, media, info, *,
               title="上节课知识点回顾", points=None):
    """回顾页：浅灰背景 + 知识点列表"""
    slide = prs.slides.add_slide(_find_blank_layout(prs))
    w, h = layout["slide_width"], layout["slide_height"]
    _add_fullscreen_bg(slide, str(media / theme["media"]["content_bg"]["file"]), w, h)

    # 标题
    _add_textbox(slide, layout["review_title"], [title], theme,
                 first_line_sz_key="page_title_sz",
                 body_sz_key="page_title_sz",
                 color_hex=theme["palette"]["title_muted"],
                 spacing_pct=100)

    # 知识点列表
    if points:
        pt_lines = [f"知识点{i}：{p}" for i, p in enumerate(points, 1)]
        _add_textbox(slide, layout["review_content"], pt_lines, theme,
                     first_line_sz_key="section_title_sz",
                     body_sz_key="body_sz",
                     color_hex=theme["palette"]["text_on_light"],
                     spacing_pct=115)
    return slide


def gen_toc(prs, theme, layout, media, info, *, sections=None):
    """目录页：浅灰背景 + 章/节概述"""
    slide = prs.slides.add_slide(_find_blank_layout(prs))
    w, h = layout["slide_width"], layout["slide_height"]
    _add_fullscreen_bg(slide, str(media / theme["media"]["content_bg"]["file"]), w, h)

    _add_textbox(slide, layout["review_title"], ["本节课主要内容概述"], theme,
                 first_line_sz_key="page_title_sz",
                 body_sz_key="page_title_sz",
                 color_hex=theme["palette"]["title_muted"],
                 spacing_pct=100)

    if sections:
        _add_textbox(slide, layout["review_content"], sections, theme,
                     first_line_sz_key="body_sz",
                     body_sz_key="body_sz",
                     color_hex=theme["palette"]["text_on_light"],
                     spacing_pct=150)
    return slide


def gen_section_header(prs, theme, layout, media, *,
                       chapter="", section_title=""):
    """章节分隔页（可选）：白色简洁"""
    slide = prs.slides.add_slide(_find_blank_layout(prs))

    lines = []
    if chapter:
        lines.append(chapter)
    if section_title:
        lines.append(section_title)

    _add_textbox(slide, layout["section_meta"], lines, theme,
                 first_line_sz_key="small_sz",
                 body_sz_key="small_sz",
                 color_hex=theme["palette"]["text_on_light"],
                 spacing_pct=150)
    return slide


def gen_reference(prs, theme, layout, media, info, *,
                  references=None):
    """引用/拓展页：浅灰背景 + 参考列表"""
    slide = prs.slides.add_slide(_find_blank_layout(prs))
    w, h = layout["slide_width"], layout["slide_height"]
    _add_fullscreen_bg(slide, str(media / theme["media"]["content_bg"]["file"]), w, h)

    _add_textbox(slide, layout["review_title"], ["本节课授课内容引用"], theme,
                 first_line_sz_key="page_title_sz",
                 body_sz_key="page_title_sz",
                 color_hex=theme["palette"]["title_muted"],
                 spacing_pct=100)

    if references:
        _add_textbox(slide, layout["review_content"], references, theme,
                     first_line_sz_key="body_sz",
                     body_sz_key="body_sz",
                     color_hex=theme["palette"]["text_on_light"],
                     spacing_pct=115)
    return slide


def gen_assignment(prs, theme, layout, media, info, *,
                   assignment_text=None):
    """作业页：浅灰背景 + 作业要求"""
    slide = prs.slides.add_slide(_find_blank_layout(prs))
    w, h = layout["slide_width"], layout["slide_height"]
    _add_fullscreen_bg(slide, str(media / theme["media"]["content_bg"]["file"]), w, h)

    _add_textbox(slide, layout["review_title"], ["作业要求"], theme,
                 first_line_sz_key="page_title_sz",
                 body_sz_key="page_title_sz",
                 color_hex=theme["palette"]["title_muted"],
                 spacing_pct=100)

    if assignment_text:
        _add_textbox(slide, layout["review_content"], assignment_text, theme,
                     first_line_sz_key="body_sz",
                     body_sz_key="body_sz",
                     color_hex=theme["palette"]["text_on_light"],
                     bold=True,
                     spacing_pct=115)
    return slide


def gen_ending(prs, theme, layout, media, info):
    """结尾页：深灰背景 + 联系方式 + Logo"""
    slide = prs.slides.add_slide(_find_blank_layout(prs))
    w, h = layout["slide_width"], layout["slide_height"]
    _add_fullscreen_bg(slide, str(media / theme["media"]["cover_bg"]["file"]), w, h)
    _add_logo(slide, str(media / theme["media"]["logo_combo"]["file"]), layout)

    lines = [
        info.get("course_name", "课程名称"),
        f"课程编号：{info.get('course_code', '')}",
        "",
        f"主讲教师：{info.get('teacher', '')}",
        f"办公邮箱：{info.get('email', '')}",
        f"办公地点：{info.get('office', '')}",
    ]
    _add_textbox(slide, layout["cover_text"], lines, theme,
                 color_hex=theme["palette"]["text_on_dark"],
                 bold=True,
                 spacing_pct=160)
    return slide


# ====================================================================
# 主题色注入
# ====================================================================

def inject_theme_colors(prs: Presentation, theme: dict):
    """替换 PPT 主题色板为 NFU 标准

    通过 slide_master → theme 关系链找到 theme part，
    再用 blob（原始 XML bytes）解析修改后回写。
    兼容 python-pptx 1.0.x 中通用 Part 无 element 属性的情况。
    """
    pal = theme["palette"]
    color_map = {
        "accent1": pal["accent1"], "accent2": pal["accent2"],
        "accent3": pal["accent3"], "accent4": pal["accent4"],
        "accent5": pal["accent5"], "accent6": pal["accent6"],
        "dk2": pal["dk2"], "lt2": pal["lt2"],
        "hlink": pal["hyperlink"], "folHlink": pal["followed_link"],
    }

    # 通过 slide master 的 rels 定位 theme part
    theme_part = None
    for sm in prs.slide_masters:
        for rId, rel in sm.part.rels.items():
            if "theme" in str(rel.target_ref):
                theme_part = rel.target_part
                break
        if theme_part:
            break

    if theme_part is None:
        raise ValueError("未找到 theme part")

    # 从 blob 读取 XML → 修改 → 回写
    root = etree.fromstring(theme_part.blob)
    cs = root.find(".//" + qn("a:clrScheme"))
    if cs is None:
        raise ValueError("theme XML 中未找到 clrScheme")

    for tag, hx in color_map.items():
        elem = cs.find(qn(f"a:{tag}"))
        if elem is not None:
            for ch in list(elem):
                elem.remove(ch)
            srgb = etree.SubElement(elem, qn("a:srgbClr"))
            srgb.set("val", hx)

    # 回写修改后的 XML
    theme_part._blob = etree.tostring(root, xml_declaration=True,
                                       encoding="UTF-8", standalone=True)


# ====================================================================
# course.yaml 读取
# ====================================================================

def load_course_yaml(path: str) -> dict:
    """从 course.yaml 提取课程信息"""
    with open(path, "r", encoding="utf-8") as f:
        data = yaml.safe_load(f)
    m = data.get("meta", data)
    return {
        "course_name": m.get("course_name", m.get("name", "")),
        "course_code": m.get("course_code", m.get("code", "")),
        "teacher": m.get("teacher", m.get("instructor", "")),
        "semester": m.get("semester", ""),
        "date": m.get("date", ""),
        "email": m.get("email", ""),
        "office": m.get("office", ""),
    }


# ====================================================================
# CLI 入口
# ====================================================================

def main():
    ap = argparse.ArgumentParser(
        description="NFU PPT 品牌注入 — 固定环节包裹模式",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
示例：
  # 基础包裹
  %(prog)s -i slides.pptx -o out.pptx --course-name "交互设计" --teacher "张三"

  # 从 course.yaml 读取
  %(prog)s -i slides.pptx -o out.pptx --course-yaml course.yaml

  # 跳过回顾页和目录页
  %(prog)s -i slides.pptx -o out.pptx --skip review,toc --course-name ...

  # 添加章节分隔页
  %(prog)s -i slides.pptx -o out.pptx --sections "第1章：基础,第2章：进阶" ...

  # 仅替换主题色
  %(prog)s -i slides.pptx -o out.pptx --theme-only
""")

    ap.add_argument("-i", "--input", required=True, help="输入 .pptx")
    ap.add_argument("-o", "--output", required=True, help="输出 .pptx")

    # 课程信息
    g = ap.add_argument_group("课程信息")
    g.add_argument("--course-name", help="课程名称")
    g.add_argument("--course-code", help="课程编号")
    g.add_argument("--teacher", help="主讲教师")
    g.add_argument("--semester", help="授课学期")
    g.add_argument("--date", help="授课时间")
    g.add_argument("--email", help="教师邮箱")
    g.add_argument("--office", help="办公地点")
    g.add_argument("--course-yaml", help="从 course.yaml 读取")

    # 功能控制
    f = ap.add_argument_group("功能控制")
    f.add_argument("--theme-only", action="store_true",
                   help="仅注入主题色板，不增减幻灯片")
    f.add_argument("--no-theme", action="store_true",
                   help="不替换主题色板")
    f.add_argument("--skip", default="",
                   help="跳过的固定环节（逗号分隔）: review,toc,reference,assignment")
    f.add_argument("--sections", default="",
                   help="章节分隔页标题（逗号分隔），例如: '第1章：基础,第2章：进阶'")

    args = ap.parse_args()

    # 加载主题
    theme = load_theme()

    # 构建课程信息
    if args.course_yaml:
        info = load_course_yaml(args.course_yaml)
    else:
        info = {k: getattr(args, k) or "" for k in
                ["course_name", "course_code", "teacher",
                 "semester", "date", "email", "office"]}
    # 命令行覆盖 yaml
    for k in ["course_name", "course_code", "teacher"]:
        v = getattr(args, k)
        if v:
            info[k] = v
    if not info.get("course_name"):
        info["course_name"] = "课程名称"

    # 加载 PPT
    prs = Presentation(args.input)
    aspect = detect_aspect(prs)
    layout_cfg = theme["layout"][aspect]
    media_dir = resolve_media_dir(aspect)
    skip = set(s.strip() for s in args.skip.split(",") if s.strip())

    print(f"📐 比例: {aspect} | 现有页数: {len(prs.slides)}")

    # ── 仅主题色模式 ─────────────────────────
    if args.theme_only:
        inject_theme_colors(prs, theme)
        prs.save(args.output)
        print(f"✅ 仅主题色注入完成 → {args.output}")
        return

    # ── 注入主题色 ──────────────────────────
    if not args.no_theme:
        try:
            inject_theme_colors(prs, theme)
            print("✅ 主题色板已替换")
        except Exception as e:
            print(f"⚠️ 主题色注入失败（非致命）: {e}")

    # ── 尾部固定环节（先添加，因为头部会移动索引）───
    original_count = len(prs.slides)

    if "reference" not in skip:
        gen_reference(prs, theme, layout_cfg, media_dir, info)
        print("✅ 引用页")

    if "assignment" not in skip:
        gen_assignment(prs, theme, layout_cfg, media_dir, info)
        print("✅ 作业页")

    gen_ending(prs, theme, layout_cfg, media_dir, info)
    print("✅ 结尾页")

    # ── 头部固定环节（从后往前插入到位置 0）────
    # 插入顺序需要反过来：最后插入的会排在最前
    insert_before = []  # 收集要插入的生成函数

    if "toc" not in skip:
        sections_list = None
        if args.sections:
            sections_list = [s.strip() for s in args.sections.split(",")]
        gen_toc(prs, theme, layout_cfg, media_dir, info, sections=sections_list)
        _move_slide(prs, len(prs.slides) - 1, 0)
        print("✅ 目录页")

    if "review" not in skip:
        gen_review(prs, theme, layout_cfg, media_dir, info)
        _move_slide(prs, len(prs.slides) - 1, 0)
        print("✅ 回顾页")

    # 封面页 — 始终生成
    gen_cover(prs, theme, layout_cfg, media_dir, info, position=0)
    print("✅ 封面页")

    # ── 可选：章节分隔页 ────────────────────
    if args.sections:
        section_titles = [s.strip() for s in args.sections.split(",")]
        # 在内容区域之前依次插入章节页（紧跟在头部固定环节之后）
        header_count = 1  # cover
        if "review" not in skip:
            header_count += 1
        if "toc" not in skip:
            header_count += 1
        for idx, st in enumerate(section_titles):
            gen_section_header(prs, theme, layout_cfg, media_dir,
                               chapter=st, section_title="")
            _move_slide(prs, len(prs.slides) - 1, header_count + idx)
        print(f"✅ {len(section_titles)} 个章节分隔页")

    # ── 保存 ──────────────────────────────
    prs.save(args.output)
    print(f"\n🎉 品牌化完成 → {args.output}")
    print(f"   总页数: {len(prs.slides)} "
          f"(原有 {original_count} + 固定环节 {len(prs.slides) - original_count})")


if __name__ == "__main__":
    main()
